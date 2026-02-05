import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation(output_path):
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Rounding Numbers"
    subtitle.text = "Year 5 Mathematics\nLearning to Estimate and Check Reasonableness"
    
    # Slide 2: The Core Rule
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Golden Rule of Rounding"
    content = slide.placeholders[1]
    content.text = "1. Look at the rounding digit.\n2. Look at the digit to its RIGHT.\n3. Apply the rule:"
    p = content.text_frame.add_paragraph()
    p.text = "5 or More: ROUND UP"
    p.font.bold = True
    p.font.size = Pt(32)
    p = content.text_frame.add_paragraph()
    p.text = "4 or Less: LET IT REST (Stay the same)"
    p.font.bold = True
    p.font.size = Pt(32)
    
    # Slide 3: Rounding to the Nearest 10
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Rounding to the Nearest 10"
    tf = slide.placeholders[1].text_frame
    tf.text = "Examples:"
    tf.add_paragraph().text = "37 -> Look at the 7. It's 5 or more. Round up to 40."
    tf.add_paragraph().text = "82 -> Look at the 2. It's 4 or less. Stay at 80."
    tf.add_paragraph().text = "1894 -> Look at the 4. Stay at 1890."
    tf.add_paragraph().text = "1895 -> Look at the 5. Round up to 1900! (Tricky!)"
    
    # Slide 4: Rounding to the Nearest 100
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Rounding to the Nearest 100"
    tf = slide.placeholders[1].text_frame
    tf.text = "Examples:"
    tf.add_paragraph().text = "742 -> Look at the 4. Round down to 700."
    tf.add_paragraph().text = "381 -> Look at the 8. Round up to 400."
    tf.add_paragraph().text = "679 -> Look at the 7. Round up to 700."
    
    # Slide 5: Rounding to the Nearest 1000
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Rounding to the Nearest 1000"
    tf = slide.placeholders[1].text_frame
    tf.text = "Examples:"
    tf.add_paragraph().text = "2790 -> Look at the 7. Round up to 3000."
    tf.add_paragraph().text = "4281 -> Look at the 2. Round down to 4000."
    tf.add_paragraph().text = "4500 -> Exactly in the middle? Rule says 5 or more, so 5000!"
    
    # Slide 6: Your Turn!
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Your Turn!"
    tf = slide.placeholders[1].text_frame
    tf.text = "Open your books to Page 30."
    tf.add_paragraph().text = "Activity 2:04 Rounding"
    tf.add_paragraph().text = "Try to finish all 6 sections. Good luck!"
    
    prs.save(output_path)
    print(f"Presentation created at {output_path}")

if __name__ == "__main__":
    output_dir = r"c:\Users\dsuth\Documents\Joshua\Signpost Math Lessons"
    os.makedirs(output_dir, exist_ok=True)
    create_presentation(os.path.join(output_dir, "Rounding_Lesson_Presentation.pptx"))
