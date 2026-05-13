import os
import sys
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')
if hasattr(sys.stderr, 'reconfigure'):
    sys.stderr.reconfigure(encoding='utf-8')
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- HELPER FUNCTIONS FOR WORD DOCUMENT EDITORIAL QUALITY ---

def set_cell_width(cell, width_inches):
    """Sets cell width natively using both python-docx property and strict DXA unit XML tags."""
    cell.width = Inches(width_inches)
    width_twips = int(width_inches * 1440)
    tcPr = cell._tc.get_or_add_tcPr()
    tcW = OxmlElement('w:tcW')
    tcW.set(qn('w:w'), str(width_twips))
    tcW.set(qn('w:type'), 'dxa')
    tcPr.append(tcW)

def set_cell_background(cell, hex_color):
    """Applies a high-contrast background shading color to a table cell."""
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)

def set_cell_margins(cell, top=100, bottom=100, left=150, right=150):
    """Sets standard internal padding (margins) for table cells."""
    tcPr = cell._tc.get_or_add_tcPr()
    tcMar = OxmlElement('w:tcMar')
    for m, val in [('top', top), ('bottom', bottom), ('left', left), ('right', right)]:
        node = OxmlElement(f'w:{m}')
        node.set(qn('w:w'), str(val))
        node.set(qn('w:type'), 'dxa')
        tcMar.append(node)
    tcPr.append(tcMar)

def create_lesson_plan(output_path):
    """Generates the premium, curriculum-aligned lesson plan Word document."""
    doc = Document()
    
    # Page Setup - Clean Margins
    sections = doc.sections
    for section in sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.0)
        section.right_margin = Inches(1.0)

    # Base Styles Configuration
    style_normal = doc.styles['Normal']
    style_normal.font.name = 'Arial'
    style_normal.font.size = Pt(11)
    style_normal.font.color.rgb = RGBColor(15, 23, 42) # Slate Dark
    style_normal.paragraph_format.line_spacing = 1.2
    style_normal.paragraph_format.space_after = Pt(6)

    # Document Title
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    title_p.paragraph_format.space_after = Pt(12)
    title_run = title_p.add_run("Mathematics Lesson Plan: Finding Missing Numbers")
    title_run.font.name = 'Arial'
    title_run.font.size = Pt(22)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(15, 23, 42) # Deep Slate

    subtitle_p = doc.add_paragraph()
    subtitle_p.paragraph_format.space_after = Pt(18)
    sub_run = subtitle_p.add_run("Mastering Inverse Operations, Balancing Equations, and the Distributive Law")
    sub_run.font.size = Pt(13)
    sub_run.font.italic = True
    sub_run.font.color.rgb = RGBColor(20, 184, 166) # Vivid Teal Accent

    # Metadata Overview Table (Dual-width, DXA compliance)
    table = doc.add_table(rows=4, cols=2)
    table.alignment = WD_ALIGN_PARAGRAPH.CENTER
    table.autofit = False

    metadata = [
        ("Target Year Level:", "Year 6 Core (Explicit mapping to Year 5 support & extension)"),
        ("Subject Area:", "Mathematics (Number & Algebra)"),
        ("Suggested Duration:", "60 Minutes"),
        ("Curriculum Alignment (AC v9):", "AC9M5A01, AC9M5A02, AC9M6A02, AC9M6N02")
    ]

    for idx, (label, val) in enumerate(metadata):
        row = table.rows[idx]
        cell_label = row.cells[0]
        cell_val = row.cells[1]
        
        set_cell_width(cell_label, 2.2)
        set_cell_width(cell_val, 4.3)
        set_cell_margins(cell_label, top=120, bottom=120, left=150, right=150)
        set_cell_margins(cell_val, top=120, bottom=120, left=150, right=150)
        
        # Zebra striping for premium feel
        bg_color = "F8FAFC" if idx % 2 == 0 else "FFFFFF"
        set_cell_background(cell_label, bg_color)
        set_cell_background(cell_val, bg_color)

        p_label = cell_label.paragraphs[0]
        p_label.paragraph_format.space_after = Pt(0)
        run_l = p_label.add_run(label)
        run_l.font.bold = True
        run_l.font.color.rgb = RGBColor(71, 85, 105)

        p_val = cell_val.paragraphs[0]
        p_val.paragraph_format.space_after = Pt(0)
        p_val.add_run(val)

    # Spacer
    doc.add_paragraph().paragraph_format.space_after = Pt(12)

    # Section Helper
    def add_section_heading(text):
        h = doc.add_paragraph()
        h.paragraph_format.space_before = Pt(16)
        h.paragraph_format.space_after = Pt(8)
        h.paragraph_format.keep_with_next = True
        run = h.add_run(text)
        run.font.name = 'Arial'
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = RGBColor(20, 184, 166) # Vivid Teal
        return h

    # 1. Learning Intention
    add_section_heading("1. Learning Intention")
    doc.add_paragraph("We are learning to find missing numbers in mathematical sentences using inverse operations, order of operations, and the distributive law to build deep algebraic reasoning and ensure the reasonableness of our calculations.")

    # 2. Success Criteria
    add_section_heading("2. Success Criteria")
    p_sc = doc.add_paragraph("I can successfully:")
    p_sc.paragraph_format.space_after = Pt(4)
    
    sc_items = [
        "Apply inverse operations to solve single-step addition, subtraction, multiplication, and division equations.",
        "Recognise square numbers and explain their connection to multiplication facts.",
        "Simplify known operations first to systematically solve multi-step numerical equations.",
        "Use the distributive law to identify equivalent grouped expressions.",
        "Check and prove the reasonableness of my solutions by substituting answers back into the original number sentence."
    ]
    for item in sc_items:
        p = doc.add_paragraph(item, style='List Bullet')
        p.paragraph_format.space_after = Pt(3)

    # 3. Introduction / Warm-up
    add_section_heading("3. Introduction & Warm-up (10 Minutes)")
    doc.add_paragraph("Commence the lesson with a rapid-fire review of core number fact families. Write an incomplete triad on the board (e.g., 6, 8, 48) and prompt students to state the two multiplication and two division facts associated with it.")
    doc.add_paragraph("Introduce the core conceptual anchor: “An equation is exactly like a balanced set of scales. The equals sign demands that the total value on the left side is identical to the total value on the right side.” Discuss why finding missing numbers is essential for real-world problem solving and estimating.")

    # 4. Direct Instruction & Guided Practice
    add_section_heading("4. Direct Instruction & Modeled Solving (20 Minutes)")
    doc.add_paragraph("Project the accompanying PowerPoint presentation to guide interactive, step-by-step discovery:")
    
    di_steps = [
        "Single-Step Inverse Operations: Guide students through Sections 1, 2, and 3 of the worksheet. Show explicitly how addition undoes subtraction, and multiplication undoes division. Model setting out clear written working.",
        "Division & Square Numbers: Address Section 4. Define squaring a number as multiplying it by itself. Reinforce that division missing numbers require careful examination of the total dividend vs. divisor.",
        "Multi-Step Equations: Model Section 5. Teach the systematic strategy of 'simplify first'. For example, in the sentence 6 + 7 + [ ] = 15, first combine 6 + 7 to get 13, reducing the sentence to 13 + [ ] = 15.",
        "Unlocking the Distributive Law: Analyze Section 6 collaboratively. Demonstrate how a factor outside a bracket distributes across terms added inside: A × (B + C) = (A × B) + (A × C). Guide students to visually match pairs rather than calculating massive totals."
    ]
    for step in di_steps:
        p = doc.add_paragraph(step, style='List Bullet')
        p.paragraph_format.space_after = Pt(4)

    # 5. Independent Practice & Differentiation
    add_section_heading("5. Independent Practice (25 Minutes)")
    doc.add_paragraph("Students open their resource books to complete the worksheet titled 'Extra Support 18: Finding missing numbers'. Circulate to monitor written recording strategies and facilitate targeted conversations.")
    
    # Differentiation Sub-table
    doc.add_paragraph().paragraph_format.space_after = Pt(4)
    diff_table = doc.add_table(rows=4, cols=2)
    diff_table.alignment = WD_ALIGN_PARAGRAPH.CENTER
    diff_table.autofit = False

    diff_headers = diff_table.rows[0]
    set_cell_width(diff_headers.cells[0], 1.5)
    set_cell_width(diff_headers.cells[1], 5.0)
    set_cell_background(diff_headers.cells[0], "0F172A") # Deep Slate Header
    set_cell_background(diff_headers.cells[1], "0F172A")
    set_cell_margins(diff_headers.cells[0], top=140, bottom=140, left=150, right=150)
    set_cell_margins(diff_headers.cells[1], top=140, bottom=140, left=150, right=150)

    p0 = diff_headers.cells[0].paragraphs[0]
    r0 = p0.add_run("Readiness Level")
    r0.font.bold = True
    r0.font.color.rgb = RGBColor(255, 255, 255)

    p1 = diff_headers.cells[1].paragraphs[0]
    r1 = p1.add_run("Targeted Strategy & Task Adjustments")
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(255, 255, 255)

    diff_rows = [
        ("Support\n(Tier 2)", "Focus primarily on Sections 1 to 4. Provide concrete manipulatives, base-ten blocks, or physical multiplication grids. Encourage verbalising the sentence before attempting to write the missing digit."),
        ("Core\n(Standard)", "Complete Sections 1 to 5 independently. Instruct students to write down intermediate steps directly above multi-step operations to maintain complete working memory focus."),
        ("Extension\n(Advanced)", "Master Section 6 independently. Challenge students to formulate three original balancing equations utilising the distributive law and verify them with peers.")
    ]

    for idx, (level, desc) in enumerate(diff_rows, start=1):
        row = diff_table.rows[idx]
        c0, c1 = row.cells[0], row.cells[1]
        set_cell_width(c0, 1.5)
        set_cell_width(c1, 5.0)
        set_cell_margins(c0, top=120, bottom=120, left=150, right=150)
        set_cell_margins(c1, top=120, bottom=120, left=150, right=150)
        
        bg = "F8FAFC" if idx % 2 != 0 else "FFFFFF"
        set_cell_background(c0, bg)
        set_cell_background(c1, bg)

        p_lvl = c0.paragraphs[0]
        p_lvl.paragraph_format.space_after = Pt(0)
        r_lvl = p_lvl.add_run(level)
        r_lvl.font.bold = True
        r_lvl.font.color.rgb = RGBColor(15, 23, 42)

        p_dsc = c1.paragraphs[0]
        p_dsc.paragraph_format.space_after = Pt(0)
        p_dsc.add_run(desc)

    # Spacer
    doc.add_paragraph().paragraph_format.space_after = Pt(12)

    # 6. Conclusion / Plenary
    add_section_heading("6. Conclusion & Consolidation (5 Minutes)")
    doc.add_paragraph("Reconvene as a whole class to examine highly rigorous or deceptive sentences, such as question 4g ([ ] squared = 9) and multi-step subtractions. Invite students to articulate the explicit thought processes used to confirm their final answers.")
    doc.add_paragraph("Conclude by reinforcing that verifying answers for reasonableness builds total mathematical confidence.")

    # Save Document
    doc.save(output_path)
    print(f"✅ Word Document Lesson Plan generated successfully at:\n   {output_path}\n")


# --- HELPER FUNCTIONS FOR POWERPOINT PRESENTATION QUALITY ---

def apply_text_frame_defaults(tf):
    """Ensures text frames have clean internal margins and proper wrapping."""
    tf.word_wrap = True
    tf.margin_left = PInches(0.1)
    tf.margin_right = PInches(0.1)
    tf.margin_top = PInches(0.1)
    tf.margin_bottom = PInches(0.1)

def add_full_bleed_background(slide, prs_width, prs_height, r, g, b):
    """Adds a guaranteed cross-platform background rectangle sent to the lowest z-order."""
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs_width, prs_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = PRGBColor(r, g, b)
    bg_shape.line.fill.background() # No border
    # Note: added first, so naturally sits at the bottom layer.
    return bg_shape

def create_presentation(output_path):
    """Generates the highly premium, custom-designed PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = PInches(10.0)
    prs.slide_height = PInches(5.625) # Standard 16:9 widescreen layout
    
    # Custom Palettes & Styling Constants
    COLOR_DARK_BG = (15, 23, 42)      # #0F172A Deep Slate
    COLOR_LIGHT_BG = (248, 250, 252)  # #F8FAFC Crisp Off-White
    COLOR_TEAL_ACCENT = (20, 184, 166) # #14B8A6 Vivid Teal
    COLOR_TEXT_DARK = (15, 23, 42)
    COLOR_TEXT_MUTED = (100, 116, 139) # #64748B Slate Grey
    COLOR_CARD_BG = (255, 255, 255)

    blank_layout = prs.slide_layouts[6] # Fully blank layout for ultimate design control

    # ---------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Sandwich Architecture Start)
    # ---------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_full_bleed_background(slide1, prs.slide_width, prs.slide_height, *COLOR_DARK_BG)

    # Accent decorative strip on left margin
    strip = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(0.25), prs.slide_height)
    strip.fill.solid()
    strip.fill.fore_color.rgb = PRGBColor(*COLOR_TEAL_ACCENT)
    strip.line.fill.background()

    # Title & Subtitle Box (Left-aligned, 0.5" minimum safe padding)
    tb_title = slide1.shapes.add_textbox(PInches(0.8), PInches(1.8), PInches(8.5), PInches(2.5))
    tf_title = tb_title.text_frame
    apply_text_frame_defaults(tf_title)
    
    p_main = tf_title.paragraphs[0]
    p_main.text = "Finding Missing Numbers"
    p_main.font.name = 'Arial'
    p_main.font.size = PPt(44)
    p_main.font.bold = True
    p_main.font.color.rgb = PRGBColor(255, 255, 255)
    p_main.space_after = PPt(12)

    p_sub = tf_title.add_paragraph()
    p_sub.text = "Mastering Inverse Operations, Balancing Scales, & the Distributive Law"
    p_sub.font.name = 'Arial'
    p_sub.font.size = PPt(20)
    p_sub.font.color.rgb = PRGBColor(*COLOR_TEAL_ACCENT)
    p_sub.space_after = PPt(8)

    p_meta = tf_title.add_paragraph()
    p_meta.text = "Year 6 Mathematics  |  Number & Algebra Framework"
    p_meta.font.name = 'Arial'
    p_meta.font.size = PPt(14)
    p_meta.font.color.rgb = PRGBColor(148, 163, 184) # Light Slate

    # ---------------------------------------------------------
    # SLIDE 2: Direct Instruction - The Golden Rules (Light Content)
    # ---------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_full_bleed_background(slide2, prs.slide_width, prs.slide_height, *COLOR_LIGHT_BG)

    # Header Title
    tb_h2 = slide2.shapes.add_textbox(PInches(0.6), PInches(0.4), PInches(8.8), PInches(0.8))
    tf_h2 = tb_h2.text_frame
    apply_text_frame_defaults(tf_h2)
    p_h2 = tf_h2.paragraphs[0]
    p_h2.text = "The Golden Rules of Balancing"
    p_h2.font.name = 'Arial'
    p_h2.font.size = PPt(28)
    p_h2.font.bold = True
    p_h2.font.color.rgb = PRGBColor(*COLOR_TEXT_DARK)

    # Content Box - Bulleted Principles
    tb_c2 = slide2.shapes.add_textbox(PInches(0.6), PInches(1.4), PInches(8.8), PInches(3.6))
    tf_c2 = tb_c2.text_frame
    apply_text_frame_defaults(tf_c2)

    rules = [
        ("The Scale Principle", "An equation is an exactly balanced scale. The values on both sides of the equals sign MUST represent the identical total."),
        ("Addition ↔ Subtraction", "They are inverse operations (exact opposites). If a missing number sentence adds a value, subtract it from the total to work backwards!"),
        ("Multiplication ↔ Division", "They form closely linked fact families. Use known multiplication facts to rapidly uncover unknown division components.")
    ]

    for idx, (title, desc) in enumerate(rules):
        p_item = tf_c2.paragraphs[0] if idx == 0 else tf_c2.add_paragraph()
        p_item.text = f"•  {title}: "
        p_item.font.name = 'Arial'
        p_item.font.size = PPt(18)
        p_item.font.bold = True
        p_item.font.color.rgb = PRGBColor(*COLOR_TEAL_ACCENT)
        p_item.space_after = PPt(4)
        if idx > 0:
            p_item.space_before = PPt(14)

        # Add description text run
        run_desc = p_item.add_run()
        run_desc.text = desc
        run_desc.font.bold = False
        run_desc.font.color.rgb = PRGBColor(*COLOR_TEXT_DARK)

    # ---------------------------------------------------------
    # SLIDE 3: Guided Practice - Single-Step Solving (Light Content)
    # ---------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_full_bleed_background(slide3, prs.slide_width, prs.slide_height, *COLOR_LIGHT_BG)

    tb_h3 = slide3.shapes.add_textbox(PInches(0.6), PInches(0.4), PInches(8.8), PInches(0.8))
    tf_h3 = tb_h3.text_frame
    apply_text_frame_defaults(tf_h3)
    p_h3 = tf_h3.paragraphs[0]
    p_h3.text = "Solving Single-Step Sentences"
    p_h3.font.name = 'Arial'
    p_h3.font.size = PPt(28)
    p_h3.font.bold = True
    p_h3.font.color.rgb = PRGBColor(*COLOR_TEXT_DARK)

    # 3 Example Cards side-by-side using shapes
    card_width = PInches(2.7)
    card_height = PInches(3.5)
    card_y = PInches(1.4)
    
    examples_s3 = [
        ("Addition", "[   ] + 8 = 15", "Inverse Operation:\n15 - 8 = 7\n\nMissing Number: 7"),
        ("Subtraction", "[   ] - 3 = 7", "Inverse Operation:\n7 + 3 = 10\n\nMissing Number: 10"),
        ("Multiplication", "8 × [   ] = 40", "Fact Family:\n40 ÷ 8 = 5\n\nMissing Number: 5")
    ]

    for idx, (lbl, eq, sol) in enumerate(examples_s3):
        card_x = PInches(0.6 + idx * 3.0)
        
        # White Card Background
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = PRGBColor(*COLOR_CARD_BG)
        card.line.color.rgb = PRGBColor(226, 232, 240) # Subtle border
        
        # Text frame inside card
        tf_card = card.text_frame
        apply_text_frame_defaults(tf_card)
        tf_card.vertical_anchor = MSO_ANCHOR.TOP
        
        p_lbl = tf_card.paragraphs[0]
        p_lbl.text = lbl.upper()
        p_lbl.font.name = 'Arial'
        p_lbl.font.size = PPt(12)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = PRGBColor(*COLOR_TEXT_MUTED)
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.space_after = PPt(14)

        p_eq = tf_card.add_paragraph()
        p_eq.text = eq
        p_eq.font.name = 'Arial'
        p_eq.font.size = PPt(22)
        p_eq.font.bold = True
        p_eq.font.color.rgb = PRGBColor(*COLOR_TEAL_ACCENT)
        p_eq.alignment = PP_ALIGN.CENTER
        p_eq.space_after = PPt(20)

        p_sol = tf_card.add_paragraph()
        p_sol.text = sol
        p_sol.font.name = 'Arial'
        p_sol.font.size = PPt(15)
        p_sol.font.color.rgb = PRGBColor(*COLOR_TEXT_DARK)
        p_sol.alignment = PP_ALIGN.CENTER

    # ---------------------------------------------------------
    # SLIDE 4: Guided Practice - Division & Squares (Light Content)
    # ---------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_full_bleed_background(slide4, prs.slide_width, prs.slide_height, *COLOR_LIGHT_BG)

    tb_h4 = slide4.shapes.add_textbox(PInches(0.6), PInches(0.4), PInches(8.8), PInches(0.8))
    tf_h4 = tb_h4.text_frame
    apply_text_frame_defaults(tf_h4)
    p_h4 = tf_h4.paragraphs[0]
    p_h4.text = "Mastering Division & Square Numbers"
    p_h4.font.name = 'Arial'
    p_h4.font.size = PPt(28)
    p_h4.font.bold = True
    p_h4.font.color.rgb = PRGBColor(*COLOR_TEXT_DARK)

    # 2 Wider Example Cards
    wcard_width = PInches(4.2)
    wcard_height = PInches(3.5)
    
    examples_s4 = [
        ("Division Insights", "24 ÷ [   ] = 3", "Think methodically:\nWhat shares 24 into groups of 3?\n24 ÷ 3 = 8\n\nMissing Number: 8"),
        ("Square Numbers", "[   ] squared = 64", "Definition:\nMultiply a number by itself.\nWhat × What = 64?\n8 × 8 = 64\n\nMissing Number: 8")
    ]

    for idx, (lbl, eq, sol) in enumerate(examples_s4):
        card_x = PInches(0.6 + idx * 4.6)
        
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, wcard_width, wcard_height)
        card.fill.solid()
        card.fill.fore_color.rgb = PRGBColor(*COLOR_CARD_BG)
        card.line.color.rgb = PRGBColor(226, 232, 240)
        
        tf_card = card.text_frame
        apply_text_frame_defaults(tf_card)
        tf_card.vertical_anchor = MSO_ANCHOR.TOP
        
        p_lbl = tf_card.paragraphs[0]
        p_lbl.text = lbl.upper()
        p_lbl.font.name = 'Arial'
        p_lbl.font.size = PPt(13)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = PRGBColor(*COLOR_TEXT_MUTED)
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.space_after = PPt(14)

        p_eq = tf_card.add_paragraph()
        p_eq.text = eq
        p_eq.font.name = 'Arial'
        p_eq.font.size = PPt(24)
        p_eq.font.bold = True
        p_eq.font.color.rgb = PRGBColor(*COLOR_TEAL_ACCENT)
        p_eq.alignment = PP_ALIGN.CENTER
        p_eq.space_after = PPt(16)

        p_sol = tf_card.add_paragraph()
        p_sol.text = sol
        p_sol.font.name = 'Arial'
        p_sol.font.size = PPt(16)
        p_sol.font.color.rgb = PRGBColor(*COLOR_TEXT_DARK)
        p_sol.alignment = PP_ALIGN.CENTER

    # ---------------------------------------------------------
    # SLIDE 5: Direct Instruction - Multi-Step Logic (Light Content)
    # ---------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_full_bleed_background(slide5, prs.slide_width, prs.slide_height, *COLOR_LIGHT_BG)

    tb_h5 = slide5.shapes.add_textbox(PInches(0.6), PInches(0.4), PInches(8.8), PInches(0.8))
    tf_h5 = tb_h5.text_frame
    apply_text_frame_defaults(tf_h5)
    p_h5 = tf_h5.paragraphs[0]
    p_h5.text = "Strategy: Simplify Known Steps First"
    p_h5.font.name = 'Arial'
    p_h5.font.size = PPt(28)
    p_h5.font.bold = True
    p_h5.font.color.rgb = PRGBColor(*COLOR_TEXT_DARK)

    # Left content box - textual strategy
    tb_str5 = slide5.shapes.add_textbox(PInches(0.6), PInches(1.4), PInches(4.0), PInches(3.5))
    tf_str5 = tb_str5.text_frame
    apply_text_frame_defaults(tf_str5)
    
    p_st1 = tf_str5.paragraphs[0]
    p_st1.text = "When facing multiple operations:"
    p_st1.font.size = PPt(18)
    p_st1.font.bold = True
    p_st1.font.color.rgb = PRGBColor(*COLOR_TEXT_DARK)
    p_st1.space_after = PPt(12)

    p_st2 = tf_str5.add_paragraph()
    p_st2.text = "1. Find the parts with complete numbers.\n2. Calculate their combined total immediately.\n3. Rewrite the simplified sentence.\n4. Solve using standard inverse logic."
    p_st2.font.size = PPt(16)
    p_st2.font.color.rgb = PRGBColor(*COLOR_TEXT_MUTED)
    p_st2.space_after = PPt(0)

    # Right Box - Step-by-step visual demonstration
    card5 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(4.8), PInches(1.4), PInches(4.6), PInches(3.5))
    card5.fill.solid()
    card5.fill.fore_color.rgb = PRGBColor(*COLOR_CARD_BG)
    card5.line.color.rgb = PRGBColor(*COLOR_TEAL_ACCENT)
    card5.line.width = PPt(2)
    
    tf_c5 = card5.text_frame
    apply_text_frame_defaults(tf_c5)
    tf_c5.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    steps_lines = [
        ("Target Sentence:", "6 + 7 + [   ] = 15", True),
        ("Step 1 (Combine):", "13 + [   ] = 15", False),
        ("Step 2 (Inverse):", "15 - 13 = 2", False),
        ("Final Answer:", "Missing Number is 2", True)
    ]
    
    for idx, (lbl, text_val, is_highlight) in enumerate(steps_lines):
        p_l = tf_c5.paragraphs[0] if idx == 0 else tf_c5.add_paragraph()
        p_l.text = f"{lbl}  "
        p_l.font.name = 'Arial'
        p_l.font.size = PPt(15)
        p_l.font.bold = False
        p_l.font.color.rgb = PRGBColor(*COLOR_TEXT_MUTED)
        p_l.alignment = PP_ALIGN.LEFT
        if idx > 0:
            p_l.space_before = PPt(8)

        run_val = p_l.add_run()
        run_val.text = text_val
        run_val.font.bold = True
        run_val.font.size = PPt(17) if is_highlight else PPt(16)
        run_val.font.color.rgb = PRGBColor(*COLOR_TEAL_ACCENT) if is_highlight else PRGBColor(*COLOR_TEXT_DARK)

    # ---------------------------------------------------------
    # SLIDE 6: Conceptual Mastery - The Distributive Law
    # ---------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_full_bleed_background(slide6, prs.slide_width, prs.slide_height, *COLOR_LIGHT_BG)

    tb_h6 = slide6.shapes.add_textbox(PInches(0.6), PInches(0.4), PInches(8.8), PInches(0.8))
    tf_h6 = tb_h6.text_frame
    apply_text_frame_defaults(tf_h6)
    p_h6 = tf_h6.paragraphs[0]
    p_h6.text = "Unlocking the Distributive Law"
    p_h6.font.name = 'Arial'
    p_h6.font.size = PPt(28)
    p_h6.font.bold = True
    p_h6.font.color.rgb = PRGBColor(*COLOR_TEXT_DARK)

    # Wide showcase card
    card6 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(0.6), PInches(1.3), PInches(8.8), PInches(3.7))
    card6.fill.solid()
    card6.fill.fore_color.rgb = PRGBColor(*COLOR_CARD_BG)
    card6.line.color.rgb = PRGBColor(226, 232, 240)

    tf_c6 = card6.text_frame
    apply_text_frame_defaults(tf_c6)
    tf_c6.vertical_anchor = MSO_ANCHOR.TOP

    p_ddef = tf_c6.paragraphs[0]
    p_ddef.text = "RULE: A factor outside brackets distributes across added terms inside."
    p_ddef.font.name = 'Arial'
    p_ddef.font.size = PPt(15)
    p_ddef.font.bold = True
    p_ddef.font.color.rgb = PRGBColor(*COLOR_TEXT_MUTED)
    p_ddef.alignment = PP_ALIGN.CENTER
    p_ddef.space_after = PPt(18)

    p_deq = tf_c6.add_paragraph()
    p_deq.text = "5 × (4 + 6)  =  (5 × [   ]) + (5 × 6)"
    p_deq.font.name = 'Arial'
    p_deq.font.size = PPt(26)
    p_deq.font.bold = True
    p_deq.font.color.rgb = PRGBColor(*COLOR_TEAL_ACCENT)
    p_deq.alignment = PP_ALIGN.CENTER
    p_deq.space_after = PPt(20)

    p_dexp = tf_c6.add_paragraph()
    p_dexp.text = "Look at the pattern matching without computing large totals:\n• The factor 5 multiplies both the 4 and the 6.\n• The right side displays (5 × 6), leaving the matching pair (5 × 4).\n\nTherefore, the missing number is precisely 4."
    p_dexp.font.name = 'Arial'
    p_dexp.font.size = PPt(16)
    p_dexp.font.color.rgb = PRGBColor(*COLOR_TEXT_DARK)
    p_dexp.alignment = PP_ALIGN.CENTER

    # ---------------------------------------------------------
    # SLIDE 7: Plenary / Independent Solving Setup (Dark Sandwich End)
    # ---------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_full_bleed_background(slide7, prs.slide_width, prs.slide_height, *COLOR_DARK_BG)

    strip7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(9.75), 0, PInches(0.25), prs.slide_height)
    strip7.fill.solid()
    strip7.fill.fore_color.rgb = PRGBColor(*COLOR_TEAL_ACCENT)
    strip7.line.fill.background()

    tb_t7 = slide7.shapes.add_textbox(PInches(0.8), PInches(1.2), PInches(8.5), PInches(3.5))
    tf_t7 = tb_t7.text_frame
    apply_text_frame_defaults(tf_t7)

    p_fmain = tf_t7.paragraphs[0]
    p_fmain.text = "Your Turn: Independent Practice"
    p_fmain.font.name = 'Arial'
    p_fmain.font.size = PPt(38)
    p_fmain.font.bold = True
    p_fmain.font.color.rgb = PRGBColor(255, 255, 255)
    p_fmain.space_after = PPt(20)

    p_fsub = tf_t7.add_paragraph()
    p_fsub.text = "1. Open your worksheet: Extra Support 18 (Finding missing numbers)\n2. Work systematically through Sections 1 to 6.\n3. Show clear working steps for multi-operation sentences.\n4. Pro-Tip: Check answers for reasonableness by substituting them back in!"
    p_fsub.font.name = 'Arial'
    p_fsub.font.size = PPt(18)
    p_fsub.font.color.rgb = PRGBColor(241, 245, 249)
    p_fsub.space_after = PPt(0)

    # Save Presentation
    prs.save(output_path)
    print(f"✅ PowerPoint Presentation generated successfully at:\n   {output_path}\n")


if __name__ == "__main__":
    # Ensure correct output directory structure
    project_dir = r"c:\Users\dsuth\Documents\Joshua\Finding_Missing_Numbers_Lesson"
    os.makedirs(project_dir, exist_ok=True)

    lp_path = os.path.join(project_dir, "Lesson_Plan_Missing_Numbers.docx")
    pptx_path = os.path.join(project_dir, "Missing_Numbers_Presentation.pptx")

    print("=======================================================")
    print("🚀 Commencing Professional Resource Compilation...")
    print("=======================================================")
    
    try:
        create_lesson_plan(lp_path)
        create_presentation(pptx_path)
        print("🎉 All deliverables successfully compiled to production standards.")
    except Exception as e:
        print(f"❌ CRITICAL ERROR during generation: {e}", file=sys.stderr)
        sys.exit(1)
